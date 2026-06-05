"""Script to generate slides/youtube_engagement_visualization_v1.pptx."""
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FIGURES = Path("outputs/figures")
OUT = Path("slides/youtube_engagement_visualization_v1.pptx")

# Palette ----------------------------------------------------------------
BG        = RGBColor(0x16, 0x21, 0x3E)  # deep navy background
SURFACE   = RGBColor(0x1E, 0x29, 0x3B)  # card surface (slightly lighter)
HDR_BG    = RGBColor(0x0F, 0x17, 0x2A)  # header strip
OFF_WHITE = RGBColor(0xE8, 0xEA, 0xF0)  # primary text
MUTED     = RGBColor(0x94, 0xA3, 0xB8)  # secondary / footnote text
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)  # accent 1 — titles, callouts
CYAN      = RGBColor(0x22, 0xD3, 0xEE)  # accent 2 — bullets, dividers

W = Inches(13.33)
H = Inches(7.5)


# Low-level helpers -------------------------------------------------------

def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()


def add_text(slide, text: str, left, top, width, height, size: int,
             bold: bool = False, color: RGBColor = None,
             align=PP_ALIGN.LEFT) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color is not None else OFF_WHITE


def add_bullets(slide, items: list[str], left, top, width, height,
                size: int = 18) -> None:
    """Bullet list with cyan dots and off-white body text."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(14)
        dot = p.add_run()
        dot.text = "•  "
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = CYAN
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = OFF_WHITE


def add_callout(slide, text: str, left, top, width, height) -> None:
    """Purple filled box used for key takeaways."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE


def add_header(slide, title: str) -> None:
    """Dark header strip with purple left accent and large title."""
    add_rect(slide, 0, 0, W, Inches(0.95), HDR_BG)
    add_rect(slide, 0, 0, Inches(0.1), Inches(0.95), PURPLE)
    add_text(slide, title,
             Inches(0.28), Inches(0.15), W - Inches(0.5), Inches(0.68),
             size=26, bold=True)


def add_image_card(slide, path: Path, card_left, card_top,
                   img_width, padding=Inches(0.1)) -> int:
    """Draw a SURFACE-colored card, then place the image inside it.
    Returns the card height so the caller can position the next element."""
    with PILImage.open(path) as img:
        iw, ih = img.size
    img_height = int(img_width * ih / iw)
    card_width = img_width + 2 * padding
    card_height = img_height + 2 * padding
    add_rect(slide, card_left, card_top, card_width, card_height, SURFACE)
    slide.shapes.add_picture(str(path),
                             card_left + padding, card_top + padding,
                             width=img_width, height=img_height)
    return card_height


# Slide builders ----------------------------------------------------------

def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_rect(slide, 0, 0, Inches(0.12), H, PURPLE)  # left accent bar

    add_text(slide, "Music Market Analysis",
             Inches(0.5), Inches(1.8), W - Inches(1.0), Inches(1.1),
             size=44, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "YouTube Engagement Visualization",
             Inches(0.5), Inches(3.15), W - Inches(1.0), Inches(0.65),
             size=22, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, "Initial Findings",
             Inches(0.5), Inches(3.75), W - Inches(1.0), Inches(0.5),
             size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # Thin cyan divider below subtitle
    add_rect(slide, Inches(4.0), Inches(4.0), W - Inches(8.0), Inches(0.04), CYAN)

    add_text(slide,
             "Combining iTunes metadata with YouTube engagement metrics "
             "to analyze music performance across platforms.",
             Inches(2.0), Inches(4.25), W - Inches(4.0), Inches(0.9),
             size=14, color=MUTED, align=PP_ALIGN.CENTER)


def bullet_slide(prs: Presentation, heading: str,
                 bullets: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_header(slide, heading)
    add_bullets(slide, bullets,
                Inches(0.8), Inches(1.2), W - Inches(1.6), H - Inches(2.2),
                size=19)
    if note:
        add_text(slide, note,
                 Inches(0.5), H - Inches(0.65), W - Inches(1.0), Inches(0.5),
                 size=11, color=MUTED)


def chart_slide(prs: Presentation, heading: str,
                images: list[Path], takeaway: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_header(slide, heading)

    padding = Inches(0.1)
    card_top = Inches(1.05)

    if len(images) == 2:
        # Tighter margins and gap to give each chart more width.
        margin = Inches(0.2)
        gap = Inches(0.2)
        img_w = (W - 2 * margin - gap - 4 * padding) // 2
        card_w = img_w + 2 * padding
        card_h = add_image_card(slide, images[0], margin, card_top, img_w, padding)
        add_image_card(slide, images[1], margin + card_w + gap, card_top, img_w, padding)
        callout_top = card_top + card_h + Inches(0.18)
        add_callout(slide, f"Takeaway:  {takeaway}",
                    Inches(0.5), callout_top, W - Inches(1.0), Inches(0.42))
    else:
        # Single chart: slightly narrowed so the takeaway has clear room below.
        img_w = Inches(9.0)
        card_w = img_w + 2 * padding
        card_h = add_image_card(slide, images[0], (W - card_w) // 2, card_top, img_w, padding)
        callout_top = card_top + card_h + Inches(0.2)
        add_callout(slide, f"Takeaway:  {takeaway}",
                    Inches(0.5), callout_top, W - Inches(1.0), Inches(0.62))


def summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_header(slide, "Summary and Next Steps")

    # Two columns: Findings (left) and Next Steps (right).
    col_w = (W - Inches(1.9)) // 2
    left1 = Inches(0.7)
    left2 = left1 + col_w + Inches(0.5)

    add_text(slide, "Findings",
             left1, Inches(1.15), col_w, Inches(0.45),
             size=15, bold=True, color=CYAN)
    add_bullets(slide,
                ["YouTube adds public engagement signals to the iTunes baseline.",
                 "Views and likes are generally aligned in the current sample; one video dominates both.",
                 "Engagement rate reveals a different ranking after normalizing."],
                left1, Inches(1.65), col_w, Inches(3.2),
                size=17)

    # Vertical divider between columns
    divider_left = left1 + col_w + Inches(0.23)
    add_rect(slide, divider_left, Inches(1.1), Inches(0.04), H - Inches(2.0), SURFACE)

    add_text(slide, "Next Steps",
             left2, Inches(1.15), col_w, Inches(0.45),
             size=15, bold=True, color=PURPLE)
    add_bullets(slide,
                ["Turn charts into slide-ready visuals.",
                 "Expand sample beyond the lofi genre.",
                 "Link YouTube engagement back to iTunes track metadata."],
                left2, Inches(1.65), col_w, Inches(3.2),
                size=17)


# Entry point -------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs)

    bullet_slide(
        prs,
        heading="Data Source Progress",
        bullets=[
            "iTunes — metadata baseline (track, artist, genre, duration)",
            "Spotify — track metadata; popularity field not in Search response",
            "YouTube — views, likes, and comments via Data API v3",
        ],
        note="All data fetched via public APIs with no scraping.",
    )

    bullet_slide(
        prs,
        heading="YouTube Metrics",
        bullets=[
            "Views — exposure and reach",
            "Likes — active positive audience response",
            "Comments — audience conversation",
            "Engagement Rate — (likes + comments) / views",
        ],
        note="Sample: 10 lofi study music videos via search.list + videos.list.",
    )

    chart_slide(
        prs,
        heading="Views and Likes by Video",
        images=[FIGURES / "youtube_views_bar.png", FIGURES / "youtube_likes_bar.png"],
        takeaway="Views and likes are generally aligned, with one video dominating both.",
    )

    chart_slide(
        prs,
        heading="Engagement Rate by Video",
        images=[FIGURES / "youtube_engagement_rate_bar.png"],
        takeaway="Engagement rate shows a different ranking after normalizing by view volume.",
    )

    summary_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
